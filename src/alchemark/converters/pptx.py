"""PPTX → Markdown converter (requires `alchemark[pptx]`)."""

from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path
from typing import Any, ClassVar

import pptx  # noqa: F401  # ensures registry skips on missing dep

from alchemark.converters.base import BaseConverter
from alchemark.core import Result
from alchemark.exceptions import ConversionError


class PptxConverter(BaseConverter):
    """Converts PowerPoint (.pptx) presentations to Markdown.

    Each slide becomes a section under a level-2 heading. When
    ``preserve_images=True`` is set on the parent ``Alchemist``, every embedded
    picture is also extracted to disk and referenced in the markdown as a
    clickable hyperlink (``[![](path)](path)``).
    """

    extensions: ClassVar[list[str]] = [".pptx"]
    name: ClassVar[str] = "pptx"

    def convert(self, path: Path) -> Result:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        try:
            prs = Presentation(str(path))
        except Exception as e:
            raise ConversionError(
                path,
                f"Could not open presentation: {e}",
                hint="The file may be corrupted or password-protected.",
            ) from e

        lines: list[str] = []
        warnings: list[str] = []

        # Set up output directory for image extraction (only if requested).
        target_dir: Path | None = None
        if self.preserve_images:
            if self.image_dir is not None:
                target_dir = Path(self.image_dir)
            else:
                target_dir = path.parent / f"{path.stem}_images"
            try:
                target_dir.mkdir(parents=True, exist_ok=True)
            except Exception as e:
                warnings.append(f"Could not create image dir '{target_dir}': {e}")
                target_dir = None

        image_count = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            title = self._get_slide_title(slide) or f"Slide {slide_idx}"
            lines.append(f"## {title}")
            lines.append("")

            # 1. Text inside text-frame shapes.
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if self._is_title_shape(shape):
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if para.level > 0:
                        lines.append("  " * (para.level - 1) + f"- {text}")
                    else:
                        lines.append(text)

            # 2. Pictures (recursively, including those nested in groups).
            if target_dir is not None:
                slide_image_idx = 0
                for shape in self._iter_picture_shapes(slide.shapes, MSO_SHAPE_TYPE):
                    image_count += 1
                    slide_image_idx += 1
                    rel_ref = self._save_picture(
                        shape,
                        target_dir,
                        path,
                        slide_idx,
                        slide_image_idx,
                        warnings,
                    )
                    if rel_ref is None:
                        continue
                    # Inline display + click-to-open hyperlink.
                    lines.append("")
                    lines.append(f"[![]({rel_ref})]({rel_ref})")

            lines.append("")

            # 3. Speaker notes.
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append(f"> **Notes:** {notes_text}")
                    lines.append("")

        metadata: dict[str, Any] = {
            "slide_count": len(prs.slides),
            "image_count": image_count,
        }

        markdown = "\n".join(lines).strip() + "\n"
        return Result(
            markdown=markdown,
            source=path,
            metadata=metadata,
            warnings=warnings,
        )

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _get_slide_title(slide: Any) -> str | None:
        if slide.shapes.title and slide.shapes.title.text:
            text = slide.shapes.title.text.strip()
            return text or None
        return None

    @staticmethod
    def _is_title_shape(shape: Any) -> bool:
        try:
            return bool(shape.is_placeholder and shape.placeholder_format.idx == 0)
        except Exception:
            return False

    @staticmethod
    def _iter_picture_shapes(shapes: Any, mso: Any) -> Iterator[Any]:
        """Yield every Picture shape, descending into groups recursively."""
        for shape in shapes:
            shape_type = getattr(shape, "shape_type", None)
            if shape_type == mso.PICTURE:
                yield shape
            elif shape_type == mso.GROUP:
                # GroupShape exposes a `.shapes` collection of children.
                inner = getattr(shape, "shapes", None)
                if inner is not None:
                    yield from PptxConverter._iter_picture_shapes(inner, mso)

    @staticmethod
    def _save_picture(
        shape: Any,
        target_dir: Path,
        source: Path,
        slide_idx: int,
        image_idx: int,
        warnings: list[str],
    ) -> str | None:
        """Persist a Picture shape to disk and return a relative path for markdown."""
        try:
            img = shape.image
            blob = img.blob
            ext = (img.ext or "png").lstrip(".")
        except Exception as e:
            warnings.append(f"Slide {slide_idx} image {image_idx}: could not read blob ({e})")
            return None

        filename = f"slide_{slide_idx:03d}_image_{image_idx:02d}.{ext}"
        out = target_dir / filename
        try:
            out.write_bytes(blob)
        except Exception as e:
            warnings.append(f"Slide {slide_idx} image {image_idx}: write failed ({e})")
            return None

        # Build a markdown-friendly path relative to the source's parent so the
        # produced .md is portable.
        try:
            rel = out.relative_to(source.parent)
            return rel.as_posix()
        except ValueError:
            return out.as_posix()
