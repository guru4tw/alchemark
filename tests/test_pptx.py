"""Tests for the PPTX converter."""

from __future__ import annotations

from pathlib import Path

import pytest

from alchemark import alchemize
from alchemark.core import Alchemist

pytest.importorskip("pptx")


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a small PPTX with title, bullets at two levels, and speaker notes."""
    from pptx import Presentation

    prs = Presentation()
    # Layout 1 = "Title and Content"
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "First Slide"

    body = slide.placeholders[1].text_frame
    body.text = "Top-level bullet"
    sub = body.add_paragraph()
    sub.text = "Nested bullet"
    sub.level = 1

    slide.notes_slide.notes_text_frame.text = "Speaker notes here."

    # Second slide with no explicit title
    prs.slides.add_slide(prs.slide_layouts[5])  # blank-ish layout

    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return out


def test_pptx_slide_title_becomes_h2(sample_pptx: Path) -> None:
    md = alchemize(sample_pptx)
    assert "## First Slide" in md


def test_pptx_bullets_preserved(sample_pptx: Path) -> None:
    md = alchemize(sample_pptx)
    assert "Top-level bullet" in md
    assert "- Nested bullet" in md


def test_pptx_speaker_notes_quoted(sample_pptx: Path) -> None:
    md = alchemize(sample_pptx)
    assert "> **Notes:**" in md
    assert "Speaker notes here." in md


def test_pptx_metadata_slide_count(sample_pptx: Path) -> None:
    a = Alchemist()
    result = a.transmute(sample_pptx)
    assert result.metadata["slide_count"] == 2


def test_pptx_slide_without_title_falls_back(sample_pptx: Path) -> None:
    md = alchemize(sample_pptx)
    # Second slide had no explicit title, so converter falls back to "Slide N"
    assert "## Slide 2" in md
