"""Tests for PPTX image extraction (Roadmap: preserve_images for PPTX)."""

from __future__ import annotations

import io
from pathlib import Path

import pytest

from alchemark import Alchemist

pytest.importorskip("pptx")
pytest.importorskip("PIL")


def _make_png_bytes(color: tuple[int, int, int] = (255, 0, 128)) -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (40, 20), color=color).save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture
def pptx_with_images(tmp_path: Path) -> Path:
    """Build a 2-slide deck where slide 1 has text + an image and slide 2 has only an image."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    seed = tmp_path / "_seed.png"
    seed.write_bytes(_make_png_bytes())

    # Slide 1: title + bullet + picture
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Slide With Picture"
    slide1.placeholders[1].text_frame.text = "Some bullet"
    slide1.shapes.add_picture(str(seed), Inches(1), Inches(2), Inches(2), Inches(1))

    # Slide 2: blank layout, only a picture (mimics image-only slides)
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide2.shapes.add_picture(str(seed), Inches(1), Inches(1), Inches(2), Inches(1))

    out = tmp_path / "deck_with_images.pptx"
    prs.save(str(out))
    return out


def test_pptx_preserve_images_default_dir(pptx_with_images: Path) -> None:
    a = Alchemist(preserve_images=True)
    result = a.transmute(pptx_with_images)

    expected_dir = pptx_with_images.parent / f"{pptx_with_images.stem}_images"
    assert expected_dir.is_dir()
    images = sorted(expected_dir.iterdir())
    # Two pictures embedded → two image files extracted.
    assert len(images) == 2
    # Names follow slide_NNN_image_MM pattern.
    assert all(p.name.startswith("slide_") for p in images)
    assert any(p.name.startswith("slide_001_") for p in images)
    assert any(p.name.startswith("slide_002_") for p in images)
    # Markdown contains clickable hyperlinks (image inside a link).
    assert "[![](" in result.markdown
    assert f"{expected_dir.name}/slide_001_image_01" in result.markdown
    assert result.metadata["image_count"] == 2
    assert result.metadata["slide_count"] == 2


def test_pptx_preserve_images_custom_dir(tmp_path: Path, pptx_with_images: Path) -> None:
    target = tmp_path / "custom_imgs"
    a = Alchemist(preserve_images=True, image_dir=target)
    result = a.transmute(pptx_with_images)

    assert target.is_dir()
    images = list(target.iterdir())
    assert len(images) == 2
    assert "[![](" in result.markdown
    assert result.metadata["image_count"] == 2


def test_pptx_preserve_images_false_does_not_extract(
    tmp_path: Path, pptx_with_images: Path
) -> None:
    a = Alchemist(preserve_images=False)
    result = a.transmute(pptx_with_images)

    assert not (pptx_with_images.parent / f"{pptx_with_images.stem}_images").exists()
    assert "[![](" not in result.markdown
    assert "![](" not in result.markdown
    assert result.metadata["image_count"] == 0


def test_pptx_image_only_slide_now_has_picture_ref(pptx_with_images: Path) -> None:
    """Slide 2 has *only* an image — without preserve_images it would be empty.
    With preserve_images=True the image reference appears under '## Slide 2'.
    """
    a = Alchemist(preserve_images=True)
    md = a.transmute(pptx_with_images).markdown
    # Look for '## Slide 2' followed by a picture ref before the next heading.
    after = md.split("## Slide 2", 1)[1]
    assert "[![](" in after
    assert "slide_002_image_01" in after
